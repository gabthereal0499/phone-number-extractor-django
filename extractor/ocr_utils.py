from PIL import Image, ImageEnhance
import pytesseract
import re
import os

# New imports for PDF and PPT
from pdf2image import convert_from_path
import pdfplumber
from pptx import Presentation

# Tesseract OCR path for Windows
pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"


# =========================================================
# IMAGE PREPROCESSING
# =========================================================

def preprocess_image(image_path):
    """
    Preprocess image for OCR:
    - Grayscale
    - Upscale
    - Boost contrast
    """
    img = Image.open(image_path).convert("L")
    img = img.resize((img.width * 2, img.height * 2))

    enhancer = ImageEnhance.Contrast(img)
    img = enhancer.enhance(2.0)

    return img


# =========================================================
# IMAGE OCR
# =========================================================

def extract_text(image_path):
    """
    Perform OCR on the preprocessed image
    """
    img = preprocess_image(image_path)

    custom_config = (
        r"--oem 3 --psm 6 "
        r"-c tessedit_char_whitelist=0123456789+"
    )

    text = pytesseract.image_to_string(img, config=custom_config)

    print("---- OCR Output (Image) ----")
    print(text)
    print("----------------------------")

    return text


# =========================================================
# PDF SUPPORT
# =========================================================

def extract_text_from_pdf(pdf_path):
    """
    Extract text from PDF:
    1. Try direct text extraction
    2. If empty → Convert pages to images → OCR
    """
    text = ""

    try:
        # Try extracting selectable text first
        with pdfplumber.open(pdf_path) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"

        # If no text found → use OCR on pages
        if not text.strip():
            print("No direct text found. Using OCR on PDF pages...")
            images = convert_from_path(pdf_path)

            for image in images:
                ocr_result = pytesseract.image_to_string(image)
                text += ocr_result + "\n"

        print("---- OCR Output (PDF) ----")
        print(text)
        print("--------------------------")

    except Exception as e:
        print(f"Error processing PDF: {e}")

    return text


# =========================================================
# PPT SUPPORT
# =========================================================

def extract_text_from_ppt(ppt_path):
    """
    Extract text from PPTX slides
    """
    text = ""

    try:
        prs = Presentation(ppt_path)

        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"

        print("---- Extracted Text (PPT) ----")
        print(text)
        print("------------------------------")

    except Exception as e:
        print(f"Error processing PPT: {e}")

    return text


# =========================================================
# TEXT CLEANING
# =========================================================

def clean_ocr_text(text):
    """
    Remove OCR noise and normalize text
    """
    text = text.replace("\n", " ").replace("\t", " ")

    # Keep only digits, + and spaces
    text = re.sub(r"[^0-9+\s]", "", text)

    # Remove multiple spaces
    text = re.sub(r"\s+", " ", text)

    return text.strip()


# =========================================================
# PHONE NUMBER EXTRACTION
# =========================================================

def extract_phone_numbers(ocr_text):
    """
    Extract ONLY valid Indian phone numbers
    """
    clean_text = clean_ocr_text(ocr_text)

    # Find possible numbers
    candidates = re.findall(r"(?:\+91\s*)?\d{10}", clean_text)

    valid_numbers = set()

    for num in candidates:
        digits = re.sub(r"\D", "", num)

        # Remove country code
        if digits.startswith("91") and len(digits) == 12:
            digits = digits[2:]

        # Strong validation
        if (
            len(digits) == 10 and
            digits[0] in "6789" and
            not all(d == digits[0] for d in digits)  # avoid 0000000000
        ):
            valid_numbers.add(digits)

    return list(valid_numbers)


# =========================================================
# TESTING
# =========================================================

if __name__ == "__main__":
    file_path = "your_file_here"

    ext = os.path.splitext(file_path)[1].lower()

    if ext in ['.jpg', '.jpeg', '.png']:
        ocr_text = extract_text(file_path)

    elif ext == '.pdf':
        ocr_text = extract_text_from_pdf(file_path)

    elif ext == '.pptx':
        ocr_text = extract_text_from_ppt(file_path)

    else:
        print("Unsupported file format")
        ocr_text = ""

    phone_numbers = extract_phone_numbers(ocr_text)

    print("---- Extracted Phone Numbers ----")
    print(phone_numbers)
    print("--------------------------------")